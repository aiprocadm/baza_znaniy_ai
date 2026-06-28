"""Tests for the document chunking utilities."""

from __future__ import annotations

import hashlib
import io
import sys
import types
from typing import List

import pytest

pytest.importorskip("openpyxl")
pytest.importorskip("pptx")

from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

if "app.ingest.service" not in sys.modules:
    service_stub = types.ModuleType("app.ingest.service")

    class IngestJob:
        def __init__(self, file_record: object | None = None, *, attempt: int = 0) -> None:
            self.file_record = file_record
            self.attempt = attempt
            self.job_record_id: int | None = None

    class IngestWorker:
        def __init__(self, service: "IngestService") -> None:
            self.service = service
            self._task = None

        def ensure_started(self) -> None:  # pragma: no cover - no-op for tests
            return None

    class IngestService:
        def __init__(self, *args: object, **kwargs: object) -> None:
            self.worker: IngestWorker | None = None

        def set_worker(self, worker: IngestWorker) -> None:
            self.worker = worker

        def ensure_background_worker(self) -> None:  # pragma: no cover - no-op
            return None

        async def stop_background_worker(self) -> None:  # pragma: no cover - no-op
            return None

    service_stub.IngestJob = IngestJob
    service_stub.IngestService = IngestService
    service_stub.IngestWorker = IngestWorker
    sys.modules["app.ingest.service"] = service_stub

import app.ingest.chunking as ingest

_chunk = ingest._chunk
_clean = ingest._clean
_get_tokenizer = ingest._get_tokenizer
parse_and_chunk = ingest.parse_and_chunk
iter_document_pages = ingest.iter_document_pages

TOKENIZER = _get_tokenizer()


def _make_text_with_tokens(count: int, token_id: int = 100) -> tuple[str, List[int]]:
    token_ids = [token_id] * count
    text = TOKENIZER.decode(token_ids)
    return text, token_ids


def _expected_windows(
    text: str, chunk: int, overlap: int, *, tokenizer=TOKENIZER
) -> List[List[int]]:
    tokens = tokenizer.encode(text)
    if not tokens:
        return []

    window = max(int(chunk), 1)
    step_overlap = max(min(int(overlap), window - 1), 0)
    step = max(window - step_overlap, 1)

    windows: List[List[int]] = []
    last_start: int | None = None
    limit = max(len(tokens) - window + 1, 0)
    for start in range(0, limit, step):
        end = min(start + window, len(tokens))
        windows.append(tokens[start:end])
        last_start = start

    tail_start = max(len(tokens) - window, 0)
    if tail_start < len(tokens) and (last_start is None or tail_start > last_start):
        windows.append(tokens[tail_start:])

    if not windows:
        windows.append(tokens[tail_start:])

    return windows


def test_chunk_handles_zero_and_single_window_sizes() -> None:
    text = "hello"

    assert _chunk(text, chunk=0, overlap=2) == list(text)
    assert _chunk(text, chunk=1, overlap=2) == list(text)


def test_chunk_rejects_non_string_input() -> None:
    with pytest.raises(TypeError):
        _chunk(123, chunk=2, overlap=0)  # type: ignore[arg-type]


def test_chunk_normalises_overlap_and_avoids_empty_chunks() -> None:
    text = "abcdefghij"
    chunks = _chunk(text, chunk=4, overlap=10)

    assert chunks[-1].endswith("hij")
    assert all(chunk for chunk in chunks)
    assert chunks.count(text[-4:]) == 1


def test_chunk_returns_empty_list_for_empty_string() -> None:
    assert _chunk("", chunk=5, overlap=2) == []


def test_chunk_returns_single_window_when_text_shorter_than_chunk() -> None:
    text = "short"
    assert _chunk(text, chunk=20, overlap=5) == [text]


def test_chunk_handles_unicode_text() -> None:
    text = "🍰 café naïve 𝛼βγ"
    pieces = _chunk(text, chunk=3, overlap=1)

    assert pieces
    assert pieces[0] == text[:3]
    assert pieces[-1] == text[-3:]
    assert all(piece for piece in pieces)


def test_chunk_progress_with_high_overlap() -> None:
    text = "abcdef"
    chunks = _chunk(text, chunk=2, overlap=5)
    assert "".join(chunks).startswith(text[:2])
    assert chunks[-1]
    assert sum(len(chunk) for chunk in chunks) >= len(text)


def test_chunk_respects_token_window_size() -> None:
    text, original_tokens = _make_text_with_tokens(1800)
    chunks = _chunk(text, chunk=900, overlap=140, encoder=TOKENIZER)
    encoded_chunks = [TOKENIZER.encode(chunk) for chunk in chunks]

    assert len(chunks) == 3
    assert encoded_chunks[0] == original_tokens[:900]
    assert encoded_chunks[1] == original_tokens[760:1660]
    assert encoded_chunks[2] == original_tokens[-900:]
    assert all(len(tokens) <= 900 for tokens in encoded_chunks)


def test_chunk_fallback_respects_token_window_slices() -> None:
    class ExpandingTokenizer:
        def encode(self, text: str) -> List[int]:
            if text == "<expand>":
                return [1]
            return [2] * len(text)

        def decode(self, tokens: List[int]) -> str:
            if tokens == [1]:
                return "x" * 1800
            return "x" * len(tokens)

    tokenizer = ExpandingTokenizer()
    chunk = 900
    overlap = 140
    text = "<expand>"

    expanded_text = tokenizer.decode(tokenizer.encode(text))
    chunks = _chunk(text, chunk=chunk, overlap=overlap, encoder=tokenizer)

    assert len(chunks) == 3
    assert chunks[0] == expanded_text[:chunk]
    assert chunks[1] == expanded_text[chunk - overlap : chunk - overlap + chunk]
    assert chunks[2] == expanded_text[-chunk:]

    encoded_chunks = [tokenizer.encode(chunk) for chunk in chunks]
    expected = _expected_windows(expanded_text, chunk, overlap, tokenizer=tokenizer)

    assert encoded_chunks == expected
    for current, nxt in zip(encoded_chunks, encoded_chunks[1:]):
        assert current[-overlap:] == nxt[:overlap]


def test_chunk_overlap_consistency() -> None:
    text, _ = _make_text_with_tokens(1500, token_id=101)
    chunks = _chunk(text, chunk=900, overlap=140, encoder=TOKENIZER)
    encoded_chunks = [TOKENIZER.encode(chunk) for chunk in chunks]

    for current, nxt in zip(encoded_chunks, encoded_chunks[1:]):
        overlap = min(140, len(current), len(nxt))
        assert current[-overlap:] == nxt[:overlap]


def test_chunk_respects_token_boundaries() -> None:
    text = "hello world " * 10
    chunk = 15
    overlap = 4

    expected = _expected_windows(text, chunk, overlap)
    chunks = _chunk(text, chunk=chunk, overlap=overlap, encoder=TOKENIZER)
    encoded_chunks = [TOKENIZER.encode(piece) for piece in chunks]

    assert encoded_chunks == expected
    assert all(len(window) <= chunk for window in encoded_chunks)


def test_chunk_overlap_adjusts_when_chunk_is_small() -> None:
    text = "a" * 12
    chunks = _chunk(text, chunk=1, overlap=5)
    encoded = [TOKENIZER.encode(piece) for piece in chunks]

    assert len(encoded) == len(TOKENIZER.encode(text))
    assert all(len(tokens) == 1 for tokens in encoded)


def test_chunk_handles_single_token_multi_char_text() -> None:
    class SingleTokenTokenizer:
        def encode(self, text: str) -> List[int]:
            return [1] if text else []

        def decode(self, tokens: List[int]) -> str:
            return "window" if tokens else ""

    tokenizer = SingleTokenTokenizer()
    text = "window"

    chunks = _chunk(text, chunk=1, overlap=0, encoder=tokenizer)

    assert chunks == list(text)


def test_chunk_small_window_char_tokenizer_fallback() -> None:
    class ExpandingTokenizer:
        def encode(self, text: str) -> List[int]:
            return [1] if text else []

        def decode(self, tokens: List[int]) -> str:
            return "expansion" if tokens else ""

    tokenizer = ExpandingTokenizer()
    chunk = 5
    overlap = 2
    text = "trigger"

    pieces = _chunk(text, chunk=chunk, overlap=overlap, encoder=tokenizer)

    expanded = tokenizer.decode(tokenizer.encode(text))
    char_tokenizer = ingest._CharTokenizer()
    expected = _expected_windows(expanded, chunk, overlap, tokenizer=char_tokenizer)
    encoded_pieces = [char_tokenizer.encode(piece) for piece in pieces]

    assert encoded_pieces == expected


def test_chunk_small_window_reencoded_branch() -> None:
    class ReencodingTokenizer:
        def encode(self, text: str) -> List[int]:
            if text == "seed":
                return [1, 2]
            if text == "ab":
                return [3, 4]
            if not text:
                return []
            return [9] * len(text)

        def decode(self, tokens: List[int]) -> str:
            if tokens in ([1, 2], [3, 4]):
                return "ab"
            if not tokens:
                return ""
            return "x" * len(tokens)

    tokenizer = ReencodingTokenizer()
    text = "seed"

    pieces = _chunk(text, chunk=5, overlap=0, encoder=tokenizer)

    assert pieces == ["ab"]
    assert tokenizer.encode(pieces[0]) == [3, 4]


def test_chunk_small_window_returns_original_tokens_when_fallback_empty() -> None:
    class FlakyStr(str):
        def __new__(cls, value: str):
            obj = super().__new__(cls, value)
            obj._first = True
            return obj

        def __bool__(self) -> bool:
            if getattr(self, "_first", False):
                object.__setattr__(self, "_first", False)
                return True
            return False

    class EmptyFallbackTokenizer:
        def encode(self, text: str) -> List[int]:
            return [7] if text == "" else [ord(ch) for ch in text]

        def decode(self, tokens: List[int]) -> str:
            return "" if tokens else ""

    text = FlakyStr("")
    tokenizer = EmptyFallbackTokenizer()

    pieces = _chunk(text, chunk=5, overlap=0, encoder=tokenizer)

    assert pieces == []


def test_chunk_with_tiny_window_uses_characters_and_handles_empty_tokens() -> None:
    assert _chunk("hello", chunk=0, overlap=2) == list("hello")

    class TruthyEmptyStr(str):
        def __new__(cls, value: str):
            obj = super().__new__(cls, value)
            return obj

        def __bool__(self) -> bool:
            return True

    empty_text = TruthyEmptyStr("")

    assert _chunk(empty_text, chunk=1, overlap=0) == []


def test_parse_and_chunk_preserves_metadata_and_tokens() -> None:
    text = "page content " * 20
    payload = text.encode("utf-8")

    chunks = parse_and_chunk("example.txt", payload)
    encoded = [TOKENIZER.encode(chunk["text"]) for chunk in chunks]
    expected = _expected_windows(_clean(text), 900, 140)

    assert encoded == expected
    assert all(chunk["file"] == "example.txt" for chunk in chunks)
    assert all(chunk["page"] == 1 for chunk in chunks)


def test_parse_and_chunk_requires_extension() -> None:
    payload = b"contents"

    assert parse_and_chunk("", payload) == []
    assert parse_and_chunk("no_extension", payload) == []


def test_parse_and_chunk_rejects_unsupported_extension() -> None:
    payload = b"contents"

    assert parse_and_chunk("example.csv", payload) == []


def _expected_sha(file: str, page: int, text: str) -> str:
    payload = f"{file}\u0000{page}\u0000{text}".encode("utf-8")
    return hashlib.sha256(payload).hexdigest()


def _make_pptx_bytes(slides: List[str]) -> bytes:
    presentation = Presentation()
    blank_index = 6 if len(presentation.slide_layouts) > 6 else 0
    layout = presentation.slide_layouts[blank_index]

    for text in slides:
        slide = presentation.slides.add_slide(layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
        textbox.text_frame.clear()
        textbox.text_frame.text = text

    stream = io.BytesIO()
    presentation.save(stream)
    stream.seek(0)
    return stream.getvalue()


def _make_xlsx_bytes(sheets: List[List[List[str]]]) -> bytes:
    workbook = Workbook()
    first = True
    for index, rows in enumerate(sheets, start=1):
        if first:
            sheet = workbook.active
            sheet.title = f"Sheet{index}"
            first = False
        else:
            sheet = workbook.create_sheet(title=f"Sheet{index}")
        for row in rows:
            sheet.append(row)
    stream = io.BytesIO()
    workbook.save(stream)
    stream.seek(0)
    return stream.getvalue()


def test_parse_and_chunk_pdf_uses_mock_tokenizer(monkeypatch: pytest.MonkeyPatch) -> None:
    texts = ["First page text", "Second page text"]

    class FakePage:
        def __init__(self, text: str) -> None:
            self._text = text

        def get_text(self, mode: str) -> str:
            assert mode == "text"
            return self._text

    class FakeDocument:
        def __init__(self, *args: object, **kwargs: object) -> None:
            self._pages = [FakePage(text) for text in texts]
            self.page_count = len(self._pages)

        def load_page(self, index: int) -> FakePage:
            return self._pages[index]

        def close(self) -> None:
            return None

    encode_calls: List[str] = []

    class RecordingTokenizer:
        def encode(self, text: str) -> List[int]:
            encode_calls.append(text)
            return [ord(ch) for ch in text]

        def decode(self, tokens: List[int]) -> str:
            return "".join(chr(token) for token in tokens)

    tokenizer = RecordingTokenizer()
    tokenizer_calls = {"count": 0}

    def fake_get_tokenizer() -> RecordingTokenizer:
        tokenizer_calls["count"] += 1
        return tokenizer

    monkeypatch.setattr(ingest, "fitz", types.SimpleNamespace(open=lambda **_: FakeDocument()))
    monkeypatch.setattr(ingest, "pdfminer_extract_pages", None)
    monkeypatch.setattr(ingest, "_get_tokenizer", fake_get_tokenizer)
    monkeypatch.setenv("RAG_CHUNK", "50")
    monkeypatch.setenv("RAG_OVERLAP", "0")

    filename = "file.pdf"
    chunks = parse_and_chunk(filename, b"binary-pdf")

    cleaned_texts = [_clean(text) for text in texts]
    expected = [
        {
            "file": filename,
            "page": index + 1,
            "sha256": _expected_sha(filename, index + 1, text),
            "text": text,
            # Provenance metadata added by the docling auto-parser flow (471163f):
            # a legacy-parsed PDF reports ocr_used=True (ext == "pdf" and backend
            # == "legacy"). This is the intended parse_and_chunk payload contract.
            "meta": {
                "document": {"file": filename, "mime_type": "application/pdf"},
                "page": {"number": index + 1},
                "chunk": {"sha256": _expected_sha(filename, index + 1, text)},
                "provenance": {
                    "parser_backend_used": "legacy",
                    "fallback_reason": None,
                    "ocr_used": True,
                },
            },
        }
        for index, text in enumerate(cleaned_texts)
    ]

    assert chunks == expected
    assert tokenizer_calls["count"] == 1
    for text in cleaned_texts:
        assert text in encode_calls


def test_iter_document_pages_pdfminer_fallback(monkeypatch: pytest.MonkeyPatch) -> None:
    class FakeLTText:
        def __init__(self, text: str) -> None:
            self._text = text

        def get_text(self) -> str:
            return self._text

    def fake_extract_pages(_stream: io.BytesIO):
        yield [FakeLTText(" First page ")]
        yield [FakeLTText("Second page")]

    monkeypatch.setattr(ingest, "fitz", None)
    monkeypatch.setattr(ingest, "pdfminer_extract_pages", fake_extract_pages)
    monkeypatch.setattr(ingest, "LTTextContainer", FakeLTText)

    pages = list(iter_document_pages("fallback.pdf", b"binary"))

    assert [number for number, _ in pages] == [1, 2]
    assert pages[0][1] == "First page"
    assert pages[1][1] == "Second page"


def test_iter_document_pages_uses_ocr_when_page_is_empty(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    class FakePage:
        def get_text(self, mode: str) -> str:
            assert mode == "text"
            return ""

        def get_pixmap(self) -> object:
            return object()

    class FakeDocument:
        def __init__(self, *args: object, **kwargs: object) -> None:
            self.page_count = 1

        def load_page(self, index: int) -> FakePage:
            assert index == 0
            return FakePage()

        def close(self) -> None:
            return None

    monkeypatch.setattr(ingest, "fitz", types.SimpleNamespace(open=lambda **_: FakeDocument()))
    monkeypatch.setattr(ingest, "_OCR_AVAILABLE", True)
    monkeypatch.setattr(ingest, "_run_ocr_on_pixmap", lambda _pixmap: "ocr text")

    pages = list(iter_document_pages("scan.pdf", b"binary"))

    assert pages == [(1, "ocr text")]


def test_parse_and_chunk_docx_uses_mock_tokenizer(monkeypatch: pytest.MonkeyPatch) -> None:
    paragraphs = ["First   paragraph", "Second paragraph"]

    class FakeParagraph:
        def __init__(self, text: str) -> None:
            self.text = text

    class FakeDocument:
        def __init__(self, _stream: object) -> None:
            self.paragraphs = [FakeParagraph(text) for text in paragraphs]

    encode_calls: List[str] = []

    class RecordingTokenizer:
        def encode(self, text: str) -> List[int]:
            encode_calls.append(text)
            return [ord(ch) for ch in text]

        def decode(self, tokens: List[int]) -> str:
            return "".join(chr(token) for token in tokens)

    tokenizer = RecordingTokenizer()
    tokenizer_calls = {"count": 0}

    def fake_get_tokenizer() -> RecordingTokenizer:
        tokenizer_calls["count"] += 1
        return tokenizer

    monkeypatch.setattr(ingest, "Document", FakeDocument)
    monkeypatch.setattr(ingest, "_get_tokenizer", fake_get_tokenizer)
    monkeypatch.setenv("RAG_CHUNK", "50")
    monkeypatch.setenv("RAG_OVERLAP", "0")

    filename = "report.docx"
    chunks = parse_and_chunk(filename, b"binary-docx")

    cleaned_text = _clean("\n".join(paragraphs))
    docx_mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    expected = [
        {
            "file": filename,
            "page": 1,
            "sha256": _expected_sha(filename, 1, cleaned_text),
            "text": cleaned_text,
            # Provenance metadata added by the docling auto-parser flow (471163f):
            # a non-PDF document reports ocr_used=False. Intended payload contract.
            "meta": {
                "document": {"file": filename, "mime_type": docx_mime},
                "page": {"number": 1},
                "chunk": {"sha256": _expected_sha(filename, 1, cleaned_text)},
                "provenance": {
                    "parser_backend_used": "legacy",
                    "fallback_reason": None,
                    "ocr_used": False,
                },
            },
        }
    ]

    assert chunks == expected
    assert tokenizer_calls["count"] == 1
    assert cleaned_text in encode_calls


def test_iter_document_pages_converts_markdown_to_plain_text() -> None:
    payload = b"# Intro\n\n**Bold** text with [link](https://example.com)"
    pages = list(iter_document_pages("notes.md", payload))

    assert len(pages) == 1
    page_number, text = pages[0]
    assert page_number == 1
    assert "Intro" in text
    assert "Bold text" in text
    assert "example.com" not in text


def test_iter_document_pages_converts_html_to_plain_text() -> None:
    payload = (
        b"<html><body><h1>Heading</h1><p>Paragraph with <strong>bold</strong>.</p></body></html>"
    )
    pages = list(iter_document_pages("page.html", payload))

    assert len(pages) == 2
    assert pages[0][0] == 1
    assert pages[0][1] == "Heading"
    assert pages[1][0] == 2
    assert pages[1][1].startswith("Paragraph with bold.")


def test_iter_document_pages_reads_pptx_slides() -> None:
    slides = ["Slide One", "Slide Two"]
    pages = list(iter_document_pages("deck.pptx", _make_pptx_bytes(slides)))

    assert len(pages) == 2
    assert [text for _, text in pages] == slides


def test_iter_document_pages_reads_xlsx_sheets() -> None:
    sheets = [["Header", "Value"], ["Row", "1"]]
    more = [["Another", "Sheet"]]
    pages = list(iter_document_pages("book.xlsx", _make_xlsx_bytes([sheets, more])))

    assert len(pages) == 2
    assert "Header Value" in pages[0][1]
    assert "Another Sheet" in pages[1][1]


def test_parse_and_chunk_accepts_markdown(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setenv("RAG_CHUNK", "200")
    monkeypatch.setenv("RAG_OVERLAP", "0")
    chunks = parse_and_chunk("notes.md", b"# Title\nContent")

    assert chunks
    assert all(chunk["file"] == "notes.md" for chunk in chunks)


def test_parse_and_chunk_accepts_pptx(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setenv("RAG_CHUNK", "200")
    monkeypatch.setenv("RAG_OVERLAP", "0")
    slides = ["Slide A", "Slide B"]
    payload = _make_pptx_bytes(slides)
    chunks = parse_and_chunk("slides.pptx", payload)

    assert len(chunks) == len(slides)
    assert [chunk["page"] for chunk in chunks] == [1, 2]


def test_parse_and_chunk_accepts_xlsx(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setenv("RAG_CHUNK", "200")
    monkeypatch.setenv("RAG_OVERLAP", "0")
    sheets = [["Header", "Value"], ["Row", "1"]]
    payload = _make_xlsx_bytes([sheets])
    chunks = parse_and_chunk("workbook.xlsx", payload)

    assert chunks
    assert chunks[0]["page"] == 1
