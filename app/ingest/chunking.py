"""Document ingestion helpers for parsing and chunking content."""

from __future__ import annotations

import hashlib
import io
import logging
import os
import re
import time
from functools import lru_cache
from dataclasses import dataclass
from typing import (
    Any,
    BinaryIO,
    Callable,
    Iterator,
    List,
    NamedTuple,
    Optional,
    Protocol,
    Union,
)

from docx import Document

try:  # pragma: no cover - PyMuPDF optional in minimal environments
    import fitz  # type: ignore[import-not-found]
except Exception:  # pragma: no cover - gracefully degrade when unavailable
    fitz = None  # type: ignore[assignment]

try:  # pragma: no cover - pdfminer optional at runtime
    from pdfminer.high_level import extract_pages as pdfminer_extract_pages
    from pdfminer.layout import LTTextContainer
except Exception:  # pragma: no cover - fallback when dependency missing
    pdfminer_extract_pages = None  # type: ignore[assignment]
    LTTextContainer = None  # type: ignore[assignment]

try:  # pragma: no cover - optional dependency for lightweight PDF parsing fallback
    from pypdf import PdfReader
except Exception:  # pragma: no cover - fallback when dependency missing
    PdfReader = None  # type: ignore[assignment]

try:  # pragma: no cover - optional dependency used for spreadsheets
    from openpyxl import load_workbook
except Exception:  # pragma: no cover - fallback for minimal environments
    load_workbook = None  # type: ignore[assignment]

try:  # pragma: no cover - optional dependency for PowerPoint parsing
    from pptx import Presentation
except Exception:  # pragma: no cover - fallback when dependency missing
    Presentation = None  # type: ignore[assignment]

try:  # pragma: no cover - optional dependency for Markdown conversion
    import markdown as markdown_lib
except Exception:  # pragma: no cover - fallback when dependency missing
    markdown_lib = None  # type: ignore[assignment]

try:  # pragma: no cover - optional dependency for OCR
    import pytesseract
except Exception:  # pragma: no cover - fallback when dependency missing
    pytesseract = None  # type: ignore[assignment]

try:  # pragma: no cover - optional dependency for OCR
    from PIL import Image
except Exception:  # pragma: no cover - fallback when dependency missing
    Image = None  # type: ignore[assignment]

import zipfile
from xml.etree import ElementTree as ET

try:  # pragma: no cover - tokenizer optional in some environments
    import tiktoken
except ImportError:  # pragma: no cover - fallback used in tests
    tiktoken = None  # type: ignore[assignment]


try:  # pragma: no cover - optional dependency shim for lightweight environments
    from app.core.config import get_settings
except ModuleNotFoundError as exc:  # pragma: no cover - triggered in minimal test envs
    from typing import NoReturn

    _IMPORT_ERROR = exc

    def _missing_get_settings() -> NoReturn:
        raise RuntimeError(
            "app.core.config.get_settings is unavailable because optional dependencies "
            "were not installed. Install the project's runtime requirements to enable "
            "OCR configuration support."
        ) from _IMPORT_ERROR

    def get_settings():  # type: ignore[override]
        return _missing_get_settings()

from app.ingest.ocr import OCRError, OCRConfig, iter_pdf_pages_with_ocr


from app.ingest.docling_backend import DoclingBackend
from app.ingest.html import html_to_plain_text, html_to_text_sections

try:  # pragma: no cover - metrics are optional during lightweight testing
    from app.observability.metrics import (
        record_docling_fallback,
        record_docling_parse,
        record_document_parse,
        record_document_ocr_pages,
    )
except ModuleNotFoundError:  # pragma: no cover - executed when metrics deps missing
    def record_document_parse(*args, **kwargs):  # type: ignore[override]
        return None

    def record_document_ocr_pages(*args, **kwargs):  # type: ignore[override]
        return None

    def record_docling_parse(*args, **kwargs):  # type: ignore[override]
        return None

    def record_docling_fallback(*args, **kwargs):  # type: ignore[override]
        return None


LOGGER = logging.getLogger(__name__)


@dataclass
class ParseResult:
    pages: list[tuple[int, str]]
    parser_backend_used: str
    fallback_reason: str | None
    ocr_used: bool
    metadata: dict[str, Any]



class _Tokenizer(Protocol):
    def encode(self, text: str) -> List[int]:
        ...

    def decode(self, tokens: List[int]) -> str:
        ...


class _CharTokenizer:
    """Fallback tokenizer operating on raw characters."""

    def encode(self, text: str) -> List[int]:
        return [ord(ch) for ch in text]

    def decode(self, tokens: List[int]) -> str:
        return "".join(chr(token) for token in tokens)


_TOKENIZER: Optional[_Tokenizer] = None

_OCR_AVAILABLE = pytesseract is not None and Image is not None


def _resolve_parser_backend(explicit_backend: str | None = None) -> str:
    """Resolve parser backend with safe fallback to legacy."""

    settings = get_settings()
    raw = (explicit_backend or getattr(settings, "document_parser_backend", "legacy") or "legacy")
    backend = str(raw).strip().lower()
    if backend not in {"legacy", "docling", "auto"}:
        LOGGER.warning("Unknown parser backend %s; falling back to legacy", backend)
        return "legacy"
    if backend == "legacy":
        return backend

    docling_enabled = bool(getattr(settings, "docling_enabled", False))
    if not docling_enabled and backend == "docling":
        LOGGER.warning("Docling backend requested but disabled; falling back to legacy")
        return "legacy"

    try:
        import docling  # type: ignore # noqa: F401
        return "docling"
    except Exception as exc:
        if backend == "docling":
            LOGGER.warning("Docling backend requested but unavailable (%s); using legacy", exc)
        return "legacy"


@lru_cache(maxsize=1)
def _ocr_config() -> OCRConfig:
    settings = get_settings()
    return OCRConfig(
        tesseract_cmd=settings.ocr_tesseract_cmd,
        dpi=int(settings.ocr_dpi),
        page_limit=settings.ocr_page_limit,
        timeout_seconds=settings.ocr_timeout_seconds,
    )


def _ensure_binary_stream(data: Union[bytes, bytearray, BinaryIO]) -> BinaryIO:
    if isinstance(data, (bytes, bytearray)):
        return io.BytesIO(data)
    if hasattr(data, "read"):
        stream = data  # type: ignore[assignment]
        try:  # pragma: no cover - not all streams are seekable
            stream.seek(0)
        except Exception:
            pass
        return stream
    raise TypeError("Unsupported binary payload type")


def _load_tiktoken(name: str) -> Optional[_Tokenizer]:
    if tiktoken is None:  # pragma: no cover - handled during tests
        return None
    try:
        return tiktoken.get_encoding(name)
    except Exception as exc:  # pragma: no cover - defensive fallback
        LOGGER.warning("Failed to load tokenizer '%s': %s", name, exc)
        try:
            return tiktoken.encoding_for_model(name)
        except Exception:  # pragma: no cover - final fallback
            return None


@lru_cache(maxsize=1)
def _default_tokenizer() -> _Tokenizer:
    use_tiktoken = os.getenv("RAG_USE_TIKTOKEN", "").strip().lower() in {
        "1",
        "true",
        "yes",
        "on",
    }
    if not use_tiktoken:
        return _CharTokenizer()

    name = os.getenv("RAG_TOKENIZER_NAME", "cl100k_base")
    tokenizer = _load_tiktoken(name)
    if tokenizer is None and name != "text-embedding-3-small":
        tokenizer = _load_tiktoken("text-embedding-3-small")
    return tokenizer or _CharTokenizer()


def _get_tokenizer() -> _Tokenizer:
    global _TOKENIZER
    if _TOKENIZER is None:
        _TOKENIZER = _default_tokenizer()
    return _TOKENIZER


def _normalise_window_size(value: int, minimum: int = 1) -> int:
    value = int(value)
    return minimum if value < minimum else value


def _normalise_overlap(chunk: int, overlap: int) -> int:
    overlap = 0 if overlap < 0 else int(overlap)
    if chunk <= 1:
        return 0
    return min(overlap, chunk - 1)


def _clean(text: str) -> str:
    """Collapse whitespace and trim the provided *text*."""

    return re.sub(r"\s+", " ", text).strip()


def _markdown_to_plain_text(markdown_content: str) -> str:
    if not markdown_content:
        return ""

    html_content = (
        markdown_lib.markdown(markdown_content)
        if markdown_lib is not None
        else _markdown_fallback_to_html(markdown_content)
    )
    return html_to_plain_text(html_content)


def _markdown_fallback_to_html(markdown_content: str) -> str:
    """Convert basic Markdown into lightweight HTML without external deps."""

    text = markdown_content
    text = re.sub(r"^#{1,6}\s*", "", text, flags=re.MULTILINE)
    text = re.sub(r"\*\*([^*]+)\*\*", r"<strong>\1</strong>", text)
    text = re.sub(r"__([^_]+)__", r"<strong>\1</strong>", text)
    text = re.sub(r"\*([^*]+)\*", r"<em>\1</em>", text)
    text = re.sub(r"_([^_]+)_", r"<em>\1</em>", text)
    text = re.sub(r"`([^`]+)`", r"<code>\1</code>", text)
    text = re.sub(r"\[(.*?)\]\([^\)]+\)", r"<a>\1</a>", text)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = text.replace("\n\n", "</p><p>")
    text = text.replace("\n", "<br/>")
    return f"<p>{text}</p>"


def _read_stream_to_bytes(stream: BinaryIO) -> bytes:
    payload = stream.read()
    if isinstance(payload, bytes):
        return payload
    if isinstance(payload, bytearray):
        return bytes(payload)
    if isinstance(payload, memoryview):  # pragma: no cover - rarely triggered
        return payload.tobytes()
    return str(payload).encode("utf-8", errors="ignore")


def _run_ocr_on_pixmap(pixmap: object) -> str:
    if not _OCR_AVAILABLE:
        return ""

    config = _ocr_config()
    try:
        png_bytes = pixmap.tobytes("png")  # type: ignore[call-arg]
    except Exception:
        try:  # pragma: no cover - secondary attempt for alternative APIs
            png_bytes = pixmap.tobytes()
        except Exception:
            return ""

    try:
        with Image.open(io.BytesIO(png_bytes)) as image:  # type: ignore[call-arg]
            if config.tesseract_cmd:
                pytesseract.pytesseract.tesseract_cmd = config.tesseract_cmd  # type: ignore[attr-defined]
            options = []
            if config.dpi:
                options.append(f"--dpi {int(config.dpi)}")
            kwargs: dict[str, object] = {}
            if config.language:
                kwargs["lang"] = config.language
            text = pytesseract.image_to_string(image, config=" ".join(options), **kwargs)
    except Exception:  # pragma: no cover - OCR failures should not break ingest
        return ""

    return _clean(text)


def _extract_text_via_ocr(page: object) -> str:
    if not _OCR_AVAILABLE or fitz is None:
        return ""

    try:
        pixmap = page.get_pixmap()  # type: ignore[attr-defined]
    except Exception:
        return ""

    return _run_ocr_on_pixmap(pixmap)


class _WindowPlan(NamedTuple):
    token_ids: List[int]
    tokenizer: _Tokenizer


def _iterate_windows(
    token_ids: List[int], *, window: int, overlap: int, tokenizer: _Tokenizer
) -> List[str]:
    total = len(token_ids)
    if total == 0:
        return []

    step_overlap = _normalise_overlap(window, overlap)
    pieces: List[str] = []
    start = 0
    while start < total:
        end = min(start + window, total)
        pieces.append(tokenizer.decode(token_ids[start:end]))
        if end >= total:
            break
        next_start = end - step_overlap
        if next_start <= start:
            next_start = start + 1
        start = next_start

    return pieces


def _handle_small_token_window(
    text: str,
    token_ids: List[int],
    *,
    window: int,
    overlap: int,
    tokenizer: _Tokenizer,
) -> Optional[_WindowPlan]:
    decoded_text = tokenizer.decode(token_ids)

    if window <= 1:
        fallback_text = decoded_text or text
        char_tokenizer = _CharTokenizer()
        char_token_ids = (
            char_tokenizer.encode(fallback_text) if fallback_text else []
        )
        return _WindowPlan(char_token_ids, char_tokenizer)

    if len(token_ids) > window:
        return None

    if decoded_text and len(decoded_text) > window:
        char_tokenizer = _CharTokenizer()
        char_token_ids = char_tokenizer.encode(decoded_text)
        return _WindowPlan(char_token_ids, char_tokenizer)

    if len(token_ids) == 1:
        fallback_text = decoded_text or text
        if fallback_text:
            char_tokenizer = _CharTokenizer()
            char_token_ids = char_tokenizer.encode(fallback_text)
            if char_token_ids:
                return _WindowPlan(char_token_ids, char_tokenizer)
        return _WindowPlan(token_ids, tokenizer)

    if decoded_text:
        try:
            reencoded = tokenizer.encode(decoded_text)
        except Exception:  # pragma: no cover - defensive fallback
            reencoded = []
        if reencoded:
            if len(reencoded) <= window and len(decoded_text) <= window:
                if reencoded == token_ids:
                    return _WindowPlan(token_ids, tokenizer)
                return _WindowPlan(reencoded, tokenizer)
            fallback_text = decoded_text
        else:
            fallback_text = decoded_text
    else:
        fallback_text = text

    if not fallback_text:
        return _WindowPlan(token_ids, tokenizer)
    char_tokenizer = _CharTokenizer()
    char_token_ids = char_tokenizer.encode(fallback_text)
    if not char_token_ids:
        return _WindowPlan(token_ids, tokenizer)

    return _WindowPlan(char_token_ids, char_tokenizer)


def _chunk(
    text: str,
    *,
    chunk: int = 900,
    overlap: int = 140,
    encoder: Optional[_Tokenizer] = None,
    token_ids: Optional[List[int]] = None,
) -> List[str]:
    """Split ``text`` into overlapping windows based on token counts."""

    if not isinstance(text, str):
        raise TypeError("text must be str")

    if not text:
        return []

    try:
        window = int(chunk)
    except Exception:
        window = 1

    window = max(window, 1)
    if window == 1:
        return list(text)

    try:
        overlap_value = int(overlap)
    except Exception:
        overlap_value = 0

    overlap_value = max(0, min(overlap_value, window - 1))
    step = max(window - overlap_value, 1)

    if encoder is None:
        encoder = _get_tokenizer()

    tokens: List[int] = []
    if encoder is not None:
        if token_ids is not None:
            tokens = list(token_ids)
        else:
            try:
                tokens = list(encoder.encode(text))
            except Exception:  # pragma: no cover - defensive fallback
                tokens = []

    def _collect_chunks(length: int, fetch: Callable[[int, int], str]) -> List[str]:
        if length <= 0:
            return []

        pieces: List[str] = []
        last_start: Optional[int] = None
        limit = max(length - window + 1, 0)
        for start in range(0, limit, step):
            end = min(start + window, length)
            piece = fetch(start, end)
            if piece:
                pieces.append(piece)
                last_start = start

        tail_start = max(length - window, 0)
        if tail_start < length and (last_start is None or tail_start > last_start):
            tail_piece = fetch(tail_start, length)
            if tail_piece:
                pieces.append(tail_piece)

        if not pieces:
            tail_piece = fetch(tail_start, length)
            if tail_piece:
                pieces.append(tail_piece)

        return pieces

    decoded_full = ""
    if tokens:
        try:
            decoded_full = encoder.decode(tokens)
        except Exception:  # pragma: no cover - fallback to character chunking
            decoded_full = ""

        if not decoded_full or (len(tokens) > window or len(decoded_full) <= window):
            try:
                return _collect_chunks(
                    len(tokens),
                    lambda start, end: encoder.decode(tokens[start:end]),
                )
            except Exception:  # pragma: no cover - fallback to character chunking
                pass

    source_text = decoded_full or text
    return _collect_chunks(
        len(source_text),
        lambda start, end: source_text[start:end],
    )


def _iter_pdf_text_pymupdf(data: bytes) -> Iterator[tuple[int, str]]:
    if fitz is None:
        raise RuntimeError("PyMuPDF is not available")

    document = fitz.open(stream=data, filetype="pdf")  # type: ignore[call-arg]
    try:
        for index in range(document.page_count):
            page = document.load_page(index)
            try:
                text = page.get_text("text") or ""
            except Exception:  # pragma: no cover - PyMuPDF edge cases
                text = ""
            cleaned = _clean(text)
            if not cleaned:
                ocr_text = _extract_text_via_ocr(page)
                status = "success" if ocr_text else "failure"
                record_document_ocr_pages(pages=1, status=status, extension="pdf")
                cleaned = ocr_text
            if cleaned:
                yield index + 1, cleaned
    finally:
        document.close()


def _iter_pdf_text_pdfminer(data: bytes) -> Iterator[tuple[int, str]]:
    if pdfminer_extract_pages is None:
        raise RuntimeError("pdfminer.six is not available")

    stream = io.BytesIO(data)
    for page_number, layout in enumerate(pdfminer_extract_pages(stream), start=1):
        fragments: List[str] = []
        for element in layout:
            if LTTextContainer is not None and isinstance(element, LTTextContainer):
                fragments.append(element.get_text())
            elif hasattr(element, "get_text"):
                fragments.append(element.get_text())
        cleaned = _clean(" ".join(part.strip() for part in fragments))
        if cleaned:
            yield page_number, cleaned


def _iter_pdf_text(data: bytes) -> Iterator[tuple[int, str]]:
    pymupdf_error: Optional[BaseException] = None
    if fitz is not None:
        try:
            found = False
            for item in _iter_pdf_text_pymupdf(data):
                found = True
                yield item
            if found:
                return
        except Exception as exc:  # pragma: no cover - PyMuPDF failure fallback
            pymupdf_error = exc
            LOGGER.warning("PyMuPDF parser failed, falling back to pdfminer: %s", exc)

    if pdfminer_extract_pages is not None:
        yield from _iter_pdf_text_pdfminer(data)
        return

    if PdfReader is not None:
        try:
            reader = PdfReader(io.BytesIO(data))
            found = False
            for page_number, page in enumerate(reader.pages, start=1):
                try:
                    text = page.extract_text() or ""
                except Exception:  # pragma: no cover - defensive fallback
                    text = ""
                cleaned = _clean(text)
                if cleaned:
                    found = True
                    yield page_number, cleaned
            if found:
                return
        except Exception as exc:  # pragma: no cover - PdfReader failure fallback
            LOGGER.warning("pypdf parser failed, no text extracted: %s", exc)

    if pymupdf_error is not None:
        raise RuntimeError("Failed to parse PDF with available backends") from pymupdf_error

    raise RuntimeError("No PDF parser backend available")


def _iter_pdf_pages(data: bytes) -> Iterator[tuple[int, str]]:
    config = _ocr_config()
    buffered: List[tuple[int, str]] = []
    any_text = False
    last_page = 0

    try:
        for page_number, text in iter_pdf_pages_with_ocr(data, config=config):
            last_page = max(last_page, page_number)
            if text:
                if not any_text and buffered:
                    for buffered_page in buffered:
                        yield buffered_page
                    buffered.clear()
                any_text = True
                yield page_number, text
            else:
                if any_text:
                    yield page_number, text
                else:
                    buffered.append((page_number, text))
    except OCRError as exc:
        LOGGER.info("OCR unavailable, falling back to text extraction: %s", exc)
        yield from _iter_pdf_text(data)
        return
    except Exception:  # pragma: no cover - unexpected OCR failure
        LOGGER.exception("Unexpected OCR failure, using text fallback")
        for page_number, text in _iter_pdf_text(data):
            if page_number > last_page:
                yield page_number, text
        return

    if not any_text:
        LOGGER.info("OCR produced no text, using text extraction fallback")
        yield from _iter_pdf_text(data)
        return

    for item in buffered:
        yield item


def _iter_docx_text(handle: BinaryIO) -> Iterator[tuple[int, str]]:
    document = Document(handle)
    buffer = io.StringIO()
    page_number = 1
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        if buffer.tell():
            buffer.write("\n")
        buffer.write(text)
        if buffer.tell() > 8000:
            cleaned = _clean(buffer.getvalue())
            if cleaned:
                yield page_number, cleaned
                page_number += 1
            buffer = io.StringIO()
    remaining = _clean(buffer.getvalue())
    if remaining:
        yield page_number, remaining


def _iter_txt_text(handle: BinaryIO) -> Iterator[tuple[int, str]]:
    wrapper = io.TextIOWrapper(handle, encoding="utf-8", errors="ignore")
    try:
        buffer = io.StringIO()
        page_number = 1
        for line in wrapper:
            stripped = line.strip()
            if not stripped:
                continue
            if buffer.tell():
                buffer.write(" ")
            buffer.write(stripped)
            if buffer.tell() > 8000:
                text = _clean(buffer.getvalue())
                if text:
                    yield page_number, text
                    page_number += 1
                buffer = io.StringIO()
        remaining = _clean(buffer.getvalue())
        if remaining:
            yield page_number, remaining
    finally:
        try:  # pragma: no cover - best effort detach
            wrapper.detach()
        except Exception:
            pass


def _iter_markdown_text(handle: BinaryIO) -> Iterator[tuple[int, str]]:
    content = handle.read()
    if isinstance(content, bytes):
        text = content.decode("utf-8", errors="ignore")
    else:  # pragma: no cover - defensive fallback for unusual streams
        text = str(content)

    cleaned = _markdown_to_plain_text(text)
    if cleaned:
        yield 1, cleaned


def _iter_html_text(handle: BinaryIO) -> Iterator[tuple[int, str]]:
    content = handle.read()
    if isinstance(content, bytes):
        text = content.decode("utf-8", errors="ignore")
    else:  # pragma: no cover - defensive fallback
        text = str(content)

    sections = html_to_text_sections(text)
    for index, section in enumerate(sections, start=1):
        yield index, section


def _iter_pptx_text(handle: BinaryIO) -> Iterator[tuple[int, str]]:
    if Presentation is not None:
        try:
            presentation = Presentation(handle)
            for slide_index, slide in enumerate(presentation.slides, start=1):
                text_parts: List[str] = []
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        for paragraph in shape.text_frame.paragraphs:  # type: ignore[union-attr]
                            runs = [run.text for run in paragraph.runs if run.text]
                            paragraph_text = " ".join(part.strip() for part in runs if part)
                            if paragraph_text:
                                text_parts.append(paragraph_text)
                    elif hasattr(shape, "text"):
                        value = str(shape.text).strip()
                        if value:
                            text_parts.append(value)
                cleaned = _clean(" ".join(text_parts))
                if cleaned:
                    yield slide_index, cleaned
            return
        except Exception:  # pragma: no cover - fallback to zip parsing
            try:
                handle.seek(0)
            except Exception:
                pass

    with zipfile.ZipFile(handle) as archive:
        slide_names = sorted(
            name
            for name in archive.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )
        for index, name in enumerate(slide_names, start=1):
            with archive.open(name) as slide:
                text_parts: List[str] = []
                for event, element in ET.iterparse(slide, events=("end",)):
                    if event == "end" and element.tag.endswith("}t"):
                        if element.text:
                            text_parts.append(element.text.strip())
                        element.clear()
                cleaned = _clean(" ".join(text_parts))
                if cleaned:
                    yield index, cleaned


def _iter_xlsx_text(handle: BinaryIO) -> Iterator[tuple[int, str]]:
    if load_workbook is None:  # pragma: no cover - dependency missing
        raise RuntimeError("openpyxl is required to parse XLSX files")

    workbook = load_workbook(handle, read_only=True, data_only=True)
    for sheet_index, sheet in enumerate(workbook.worksheets, start=1):
        buffer = io.StringIO()
        for row in sheet.iter_rows(values_only=True):
            values = [str(cell).strip() for cell in row if cell not in (None, "")]
            if not values:
                continue
            if buffer.tell():
                buffer.write("\n")
            buffer.write(" ".join(values))
        text = _clean(buffer.getvalue())
        if text:
            yield sheet_index, text


def _hash_chunk(file: str, page: int, text: str) -> str:
    payload = f"{file}\u0000{page}\u0000{text}".encode("utf-8")
    return hashlib.sha256(payload).hexdigest()


def parse_document(filename: str, data: Union[bytes, bytearray, BinaryIO]) -> ParseResult:
    name = (filename or "").strip()
    if not name or "." not in name:
        return ParseResult([], "legacy", None, False, {"document": {}, "pages": [], "chunks": []})

    ext = name.rsplit(".", 1)[-1].lower()
    mime_by_ext = {"pdf": "application/pdf", "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation", "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "txt": "text/plain", "md": "text/markdown", "markdown": "text/markdown", "html": "text/html", "htm": "text/html"}
    mime = mime_by_ext.get(ext, "application/octet-stream")
    stream = _ensure_binary_stream(data)
    raw_bytes = _read_stream_to_bytes(stream)

    backend = _resolve_parser_backend()
    LOGGER.info("Document parser backend selected", extra={"backend_selected": backend, "mime": mime, "tenant": os.getenv("TENANT", "unknown"), "document_id": os.getenv("DOCUMENT_ID", "unknown"), "document_name": name})
    parser_backend_used = "legacy"
    fallback_reason = None
    pages: list[tuple[int, str]] = []

    if backend in {"docling", "auto"} and mime in DoclingBackend.SUPPORTED_MIME:
        docling_start = time.perf_counter()
        try:
            pages = DoclingBackend().parse(name, bytes(raw_bytes))
            parser_backend_used = "docling"
            record_docling_parse("success", time.perf_counter() - docling_start)
        except Exception as exc:
            fallback_reason = str(exc)
            record_docling_parse("error", time.perf_counter() - docling_start)
            record_docling_fallback(fallback_reason)
            LOGGER.warning("Docling parse failed for %s: %s. Fallback to legacy.", name, exc)

    if not pages:
        def _buffer() -> BinaryIO:
            return io.BytesIO(raw_bytes)
        if ext == "pdf":
            pages = list(_iter_pdf_pages(bytes(raw_bytes)))
        elif ext == "docx":
            pages = list(_iter_docx_text(_buffer()))
        elif ext == "txt":
            pages = list(_iter_txt_text(_buffer()))
        elif ext in {"md", "markdown"}:
            pages = list(_iter_markdown_text(_buffer()))
        elif ext in {"html", "htm"}:
            pages = list(_iter_html_text(_buffer()))
        elif ext == "pptx":
            pages = list(_iter_pptx_text(_buffer()))
        elif ext == "xlsx":
            pages = list(_iter_xlsx_text(_buffer()))
        parser_backend_used = "legacy" if parser_backend_used != "docling" else parser_backend_used

    metadata = {"document": {"file": name, "mime_type": mime}, "pages": [{"page": p, "text_length": len(t)} for p,t in pages], "chunks": []}
    ocr_used = ext == "pdf" and parser_backend_used == "legacy"
    LOGGER.info(
        "document_parse_backend_selected",
        extra={
            "backend_selected": backend,
            "parser_backend_used": parser_backend_used,
            "fallback_reason": fallback_reason,
            "mime": mime,
            "tenant": os.getenv("TENANT_ID", "default"),
            "document_id": _hash_chunk(name, 0, name)[:16],
        },
    )
    return ParseResult(pages, parser_backend_used, fallback_reason, ocr_used, metadata)


def iter_document_pages(filename: str, data: Union[bytes, bytearray, BinaryIO]) -> Iterator[tuple[int, str]]:
    return iter(parse_document(filename, data).pages)


def parse_and_chunk(filename: str, data: Union[bytes, bytearray, BinaryIO]) -> List[dict[str, object]]:
    """Parse ``data`` according to ``filename`` extension and chunk the text."""

    name = (filename or "").strip()
    extension = name.rsplit(".", 1)[-1].lower() if "." in name else ""

    start = time.perf_counter()
    status = "success"
    chunks: List[dict[str, object]] = []

    try:
        chunk_size = _normalise_window_size(int(os.getenv("RAG_CHUNK", "900")))
        overlap = _normalise_overlap(chunk_size, int(os.getenv("RAG_OVERLAP", "140")))
        tokenizer = _get_tokenizer()

        parse_result = parse_document(name, data)
        for page_number, page_text in parse_result.pages:
            if not page_text:
                continue
            for piece in _chunk(page_text, chunk=chunk_size, overlap=overlap, encoder=tokenizer):
                sha = _hash_chunk(name, page_number, piece)
                chunks.append(
                    {
                        "file": name,
                        "page": page_number,
                        "sha256": sha,
                        "text": piece,
                        "meta": {
                            "document": parse_result.metadata.get("document", {}),
                            "page": {"number": page_number},
                            "chunk": {"sha256": sha},
                            "provenance": {
                                "parser_backend_used": parse_result.parser_backend_used,
                                "fallback_reason": parse_result.fallback_reason,
                                "ocr_used": parse_result.ocr_used,
                            },
                        },
                    }
                )
        return chunks
    except Exception:
        status = "error"
        raise
    finally:
        duration = time.perf_counter() - start
        record_document_parse(extension, status, len(chunks), duration)


__all__ = [
    "_chunk",
    "_clean",
    "_get_tokenizer",
    "iter_document_pages",
    "parse_document",
    "parse_and_chunk",
]
